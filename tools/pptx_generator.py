from pptx import Presentation
from pptx.util import Inches
from tools.chart_builder import build_chart
from tools.layout_engine import detect_slide_type, choose_layout, SlideType


def get_body_placeholder(slide):
    for shape in slide.placeholders:
        # skip title placeholder
        if shape == slide.shapes.title:
            continue

        # find text-capable placeholder
        if shape.has_text_frame:
            return shape

    return None

def clear_unused_text_placeholders(slide):
    for shape in slide.placeholders:
        if shape != slide.shapes.title and shape.has_text_frame:
            shape.text_frame.clear()

def remove_body_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            phf = shape.placeholder_format

            # skip title placeholder
            if phf.type == 1:  # TITLE placeholder
                continue

            # remove all other placeholders
            slide.shapes._spTree.remove(shape._element)
def generate_ppt(deck_spec, output_path="output.pptx"):
    prs = Presentation()



    for slide_spec in deck_spec.slides:
        slide_type = detect_slide_type(slide_spec)
        slide_layout = choose_layout(prs, slide_type)
        slide = prs.slides.add_slide(slide_layout)
        clear_unused_text_placeholders(slide)
        if slide_type in [SlideType.CHART, SlideType.TITLE_ONLY]:
            remove_body_placeholders(slide)

        # if slide_spec.chart:
        #     slide_layout = prs.slide_layouts[5]  # Title Only layout
        # else:
        #     slide_layout = prs.slide_layouts[1]  # Title + Content
        #
        # slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_spec.title

        content = get_body_placeholder(slide)

        if slide_type == SlideType.BULLET and content:
            tf = content.text_frame
            tf.clear()
            tf.text = slide_spec.bullets[0]

            for bullet in slide_spec.bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        if slide_spec.chart:
            x = Inches(0.8)
            y = Inches(1.5)
            cx = Inches(8.5)
            cy = Inches(4.5)
            build_chart(slide, slide_spec.chart, x, y, cx, cy)

        elif slide_type == SlideType.MIXED:
            placeholders = [
                p for p in slide.placeholders
                if p != slide.shapes.title and p.has_text_frame
            ]

            if len(placeholders) >= 2:
                left = placeholders[0]
                right = placeholders[1]

                # Bullets on left
                tf = left.text_frame
                tf.clear()
                tf.text = slide_spec.bullets[0]
                for bullet in slide_spec.bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

                # Chart inside right placeholder bounds
                x = right.left
                y = right.top
                cx = right.width
                cy = right.height

                build_chart(slide, slide_spec.chart, x, y, cx, cy)

        # content = get_body_placeholder(slide)
        #
        # if slide_spec.bullets and content:
        #     tf = content.text_frame
        #     tf.clear()
        #
        #     tf.text = slide_spec.bullets[0]
        #
        #     for bullet in slide_spec.bullets[1:]:
        #         p = tf.add_paragraph()
        #         p.text = bullet
        #         p.level = 0


        # if slide_spec.bullets:
        #     content.text = slide_spec.bullets[0]
        #     for bullet in slide_spec.bullets[1:]:
        #         p = content.text_frame.add_paragraph()
        #         p.text = bullet
        #         p.level = 1



    prs.save(output_path)
    return output_path