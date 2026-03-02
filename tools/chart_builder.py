from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def build_chart(slide, chart_spec, x, y, cx, cy):
    chart_spec.validate_lengths()

    chart_data = ChartData()
    chart_data.categories = chart_spec.categories
    chart_data.add_series(chart_spec.title, chart_spec.values)

    chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
    }

    chart_type = chart_type_map.get(chart_spec.type.lower())
    if not chart_type:
        raise ValueError(f"Unsupported chart type: {chart_spec.type}")

    # --- SMART LAYOUT (prevents overlap) ---
    # x = Inches(0.8)
    # y = Inches(1.5)
    # cx = Inches(8.5)
    # cy = Inches(4.5)

    chart = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data
    ).chart

    # --- LEGEND ---
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # --- TITLE ---
    # chart.has_title = True
    # chart.chart_title.text_frame.text = chart_spec.title
    chart.has_title = False

    # --- AXIS FIXES ---
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)

    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(12)

    # prevent auto weird scaling
    value_axis.has_major_gridlines = True

    # --- DATA LABELS ---
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(11)

    # --- AUTO ROTATE LABELS IF MANY ---
    if len(chart_spec.categories) > 5:
        category_axis.tick_labels.rotation = -45

    return chart